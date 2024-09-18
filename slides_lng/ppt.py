import requests
from tempfile import NamedTemporaryFile
from pptx import Presentation
from pptx.util import Inches
from typing import Sequence, Mapping

def download_file(url) -> str:
    with requests.get(url, stream=True) as r:
        r.raise_for_status()
        with NamedTemporaryFile(suffix='.png', delete=False) as f:
            for chunk in r.iter_content(chunk_size=8192): 
                f.write(chunk)
            f.close()
        return f.name

def _add_slide(ppt: Presentation, slide: Mapping):
    bullet_slide_layout = ppt.slide_layouts[1]
    ppt_slide = ppt.slides.add_slide(bullet_slide_layout)
    shapes = ppt_slide.shapes

    # Title
    title_shape = shapes.title
    title_shape.text = slide.get("title_text", "")

    # Body
    if "text" in slide:
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        for bullet in slide.get("text", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

            if "p1" in slide:
                p = tf.add_paragraph()
                p.text = slide.get("p1")
                p.level = 1

    if "img_path" in slide:
        cur_left = 6
        for img_path in slide.get("img_path", []):
            top = Inches(2)
            left = Inches(cur_left)
            height = Inches(4)
            if img_path.startswith('http://') or img_path.startswith('https://'):
                try:
                    tmp_img_path = download_file(img_path)
                    ppt_slide.shapes.add_picture(tmp_img_path, left, top, height=height)
                    cur_left += 1
                except requests.HTTPError as ex:
                    print(f'Cannot download image {img_path}')
                    pass

def _add_title_slide(ppt: Presentation, slide: Mapping):
    title_slide_layout = ppt.slide_layouts[0]
    ppt_slide = ppt.slides.add_slide(title_slide_layout)
    title = ppt_slide.shapes.title
    subtitle = ppt_slide.placeholders[1]
    if "title_text" in slide:
        title.text = slide.get("title_text")
    if "subtitle_text" in slide:
        subtitle.text = slide.get("subtitle_text")


def create_presentation(slides: Sequence, file_name: str | None = None) -> str:

    if not file_name:
        file_name = slides[0].get("title_text").lower().replace(",", "").replace(" ", "-") + '.ppt'
        file_name = './out/' + file_name

    ppt = Presentation()
    _add_title_slide(ppt, slides[0])

    for slide in slides[1:]:
        _add_slide(ppt, slide)

    ppt.save(file_name)
    return file_name